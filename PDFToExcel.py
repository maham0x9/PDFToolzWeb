import os
import io
import fitz  # PyMuPDF
import pdfplumber # For Excel
import pandas as pd # For Excel
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
from werkzeug.utils import secure_filename

app = Flask(__name__)
CORS(app)

# Configure folders
UPLOAD_FOLDER = 'uploads'
DOWNLOAD_FOLDER = 'downloads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DOWNLOAD_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['DOWNLOAD_FOLDER'] = DOWNLOAD_FOLDER

# --- HELPER: PDF TO PPTX ---
def convert_pdf_to_pptx_logic(pdf_path, pptx_path):
    prs = Presentation()
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_stream = io.BytesIO(pix.tobytes("png"))
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    doc.close()
    prs.save(pptx_path)

# --- HELPER: PDF TO EXCEL ---
def convert_pdf_to_excel_logic(pdf_path, excel_path):
    with pdfplumber.open(pdf_path) as pdf:
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            tables_found = False
            for i, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                if tables:
                    tables_found = True
                    for j, table in enumerate(tables):
                        # Convert list of lists to DataFrame
                        df = pd.DataFrame(table)
                        # Clean up text (remove newlines within cells)
                        df = df.replace(r'\n', ' ', regex=True)
                        
                        # Create unique sheet name
                        sheet_name = f"Page{i+1}_Table{j+1}"
                        # Save to excel
                        df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
            
            # If no tables found, create a blank sheet with a message
            if not tables_found:
                df = pd.DataFrame(["No detected tables in this PDF."])
                df.to_excel(writer, sheet_name="Info", index=False, header=False)

# --- ROUTES ---

@app.route('/convert-to-word', methods=['POST'])
def convert_to_word():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        word_filename = filename.rsplit('.', 1)[0] + '.docx'
        word_path = os.path.join(app.config['DOWNLOAD_FOLDER'], word_filename)
        try:
            cv = Converter(pdf_path)
            cv.convert(word_path, start=0, end=None)
            cv.close()
            return jsonify({'message': 'Success', 'download_url': f'/download/{word_filename}'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500

@app.route('/convert-to-ppt', methods=['POST'])
def convert_to_ppt():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        ppt_filename = filename.rsplit('.', 1)[0] + '.pptx'
        ppt_path = os.path.join(app.config['DOWNLOAD_FOLDER'], ppt_filename)
        try:
            convert_pdf_to_pptx_logic(pdf_path, ppt_path)
            return jsonify({'message': 'Success', 'download_url': f'/download/{ppt_filename}'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500

@app.route('/convert-to-excel', methods=['POST'])
def convert_to_excel():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)

        excel_filename = filename.rsplit('.', 1)[0] + '.xlsx'
        excel_path = os.path.join(app.config['DOWNLOAD_FOLDER'], excel_filename)

        try:
            convert_pdf_to_excel_logic(pdf_path, excel_path)
            return jsonify({
                'message': 'Success', 
                'download_url': f'/download/{excel_filename}'
            })
        except Exception as e:
            print(e)
            return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    return send_file(os.path.join(app.config['DOWNLOAD_FOLDER'], filename), as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True, port=5000)