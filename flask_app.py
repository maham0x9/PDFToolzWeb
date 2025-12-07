import os
import io
import zipfile
import fitz  # PyMuPDF
import pdfplumber
import pandas as pd
from PIL import Image
from flask import Flask, request, send_file, jsonify, render_template, url_for
from flask_cors import CORS
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
from werkzeug.utils import secure_filename

app = Flask(__name__, template_folder='../templates', static_folder='../static')
CORS(app)

# --- CONFIGURATION ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
# Go up one level to reach the project root for uploads/downloads
PROJECT_ROOT = os.path.dirname(BASE_DIR)

UPLOAD_FOLDER = os.path.join(PROJECT_ROOT, 'uploads')
DOWNLOAD_FOLDER = os.path.join(PROJECT_ROOT, 'downloads')

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DOWNLOAD_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['DOWNLOAD_FOLDER'] = DOWNLOAD_FOLDER

# --- HELPER FUNCTIONS ---

def get_size_format(b, factor=1024, suffix="B"):
    for unit in ["", "K", "M", "G", "T"]:
        if b < factor:
            return f"{b:.2f}{unit}{suffix}"
        b /= factor
    return f"{b:.2f}Y{suffix}"

def compress_images_in_pdf(doc, quality=50, max_width=1024):
    img_xrefs = set()
    for page_num in range(len(doc)):
        page = doc[page_num]
        images = page.get_images()
        for img in images:
            xref = img[0]
            if xref in img_xrefs: continue
            img_xrefs.add(xref)
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.width < 100 or pix.height < 100: continue
                if pix.n - pix.alpha > 3: pix = fitz.Pixmap(fitz.csRGB, pix, 0)
                img_data = pix.tobytes()
                pil_img = Image.open(io.BytesIO(img_data))
                if pil_img.width > max_width:
                    ratio = max_width / float(pil_img.width)
                    new_height = int(float(pil_img.height) * ratio)
                    pil_img = pil_img.resize((max_width, new_height), Image.Resampling.LANCZOS)
                buffer = io.BytesIO()
                pil_img.save(buffer, format="JPEG", quality=quality, optimize=True)
                doc.update_stream(xref, buffer.getvalue())
            except Exception as e:
                print(f"Skipping image {xref}: {e}")

def convert_pdf_to_pptx_logic(pdf_path, pptx_path):
    prs = Presentation()
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_stream = io.BytesIO(pix.tobytes("png"))
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    doc.close()
    prs.save(pptx_path)

def convert_pdf_to_excel_logic(pdf_path, excel_path):
    with pdfplumber.open(pdf_path) as pdf:
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            tables_found = False
            for i, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                if tables:
                    tables_found = True
                    for j, table in enumerate(tables):
                        df = pd.DataFrame(table)
                        df = df.replace(r'\n', ' ', regex=True)
                        sheet_name = f"Page{i+1}_Table{j+1}"
                        df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
            if not tables_found:
                df = pd.DataFrame(["No detected tables in this PDF."])
                df.to_excel(writer, sheet_name="Info", index=False, header=False)

def parse_page_string(order_str, total_pages):
    selected_pages = []
    if not order_str: return list(range(total_pages))
    parts = order_str.split(',')
    for part in parts:
        part = part.strip()
        if not part: continue
        try:
            if '-' in part:
                start, end = map(int, part.split('-'))
                step = 1 if start <= end else -1
                r_start = max(1, start) - 1
                r_end = min(total_pages, end) - 1
                if step == 1: selected_pages.extend(range(r_start, r_end + 1))
                else: selected_pages.extend(range(r_start, r_end - 1, -1))
            else:
                p = int(part) - 1
                if 0 <= p < total_pages: selected_pages.append(p)
        except ValueError: pass
    return selected_pages if selected_pages else list(range(total_pages))

# --- FRONTEND ROUTES (Serving HTML) ---

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/tool/compress')
def view_compress():
    return render_template('CompressPDF.html')

@app.route('/tool/merge')
def view_merge():
    return render_template('MergePDF.html')

@app.route('/tool/organize')
def view_organize():
    return render_template('OrganizePDF.html')

@app.route('/tool/excel')
def view_excel():
    return render_template('PDFtoExcel.html')

@app.route('/tool/ppt')
def view_ppt():
    return render_template('PDFtoPPT.html')

# --- API ROUTES (Processing Logic) ---

@app.route('/compress-pdf', methods=['POST'])
def compress_pdf():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400

    filename = secure_filename(file.filename)
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(pdf_path)
    
    original_size = os.path.getsize(pdf_path)
    level = request.form.get('level', 'recommended') 
    base_name = filename.rsplit('.', 1)[0]
    output_filename = f"{base_name}_compressed.pdf"
    output_path = os.path.join(app.config['DOWNLOAD_FOLDER'], output_filename)

    try:
        doc = fitz.open(pdf_path)
        if level == 'extreme':
            compress_images_in_pdf(doc, quality=30, max_width=800)
            doc.save(output_path, garbage=4, deflate=True, clean=True)
        elif level == 'recommended':
            compress_images_in_pdf(doc, quality=60, max_width=1600)
            doc.save(output_path, garbage=4, deflate=True)
        else:
            doc.save(output_path, garbage=3, deflate=True)
        doc.close()

        new_size = os.path.getsize(output_path)
        if new_size >= original_size:
            doc = fitz.open(pdf_path)
            doc.save(output_path)
            doc.close()
            new_size = original_size

        return jsonify({
            'message': 'Compression successful',
            'download_url': f'/download/{output_filename}',
            'size_comparison': f"{get_size_format(original_size)} ➔ {get_size_format(new_size)}"
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/merge-pdfs', methods=['POST'])
def merge_pdfs():
    uploaded_files = request.files.getlist('files')
    if not uploaded_files or uploaded_files[0].filename == '': return jsonify({'error': 'No files selected'}), 400
    try:
        result_doc = fitz.open()
        for file in uploaded_files:
            file_stream = file.read()
            src_doc = fitz.open("pdf", file_stream)
            result_doc.insert_pdf(src_doc)
            src_doc.close()
        
        first_name = secure_filename(uploaded_files[0].filename).rsplit('.', 1)[0]
        output_filename = f"Merged_{first_name}_and_others.pdf"
        output_path = os.path.join(app.config['DOWNLOAD_FOLDER'], output_filename)
        result_doc.save(output_path)
        result_doc.close()
        
        return jsonify({'message': 'Merge successful', 'download_url': f'/download/{output_filename}'})
    except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/organize-pdf', methods=['POST'])
def organize_pdf():
    if 'file' not in request.files: return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    filename = secure_filename(file.filename)
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(pdf_path)
    page_order = request.form.get('page_order', '')
    output_filename = f"organized_{filename}"
    output_path = os.path.join(app.config['DOWNLOAD_FOLDER'], output_filename)
    try:
        doc = fitz.open(pdf_path)
        indices = parse_page_string(page_order, len(doc))
        doc.select(indices)
        doc.save(output_path)
        doc.close()
        return jsonify({'message': 'Success', 'download_url': f'/download/{output_filename}'})
    except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/convert-to-excel', methods=['POST'])
def convert_to_excel():
    if 'file' not in request.files: return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    filename = secure_filename(file.filename)
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(pdf_path)
    excel_filename = filename.rsplit('.', 1)[0] + '.xlsx'
    excel_path = os.path.join(app.config['DOWNLOAD_FOLDER'], excel_filename)
    try:
        convert_pdf_to_excel_logic(pdf_path, excel_path)
        return jsonify({'message': 'Success', 'download_url': f'/download/{excel_filename}'})
    except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/convert-to-ppt', methods=['POST'])
def convert_to_ppt():
    if 'file' not in request.files: return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    filename = secure_filename(file.filename)
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(pdf_path)
    ppt_filename = filename.rsplit('.', 1)[0] + '.pptx'
    ppt_path = os.path.join(app.config['DOWNLOAD_FOLDER'], ppt_filename)
    try:
        convert_pdf_to_pptx_logic(pdf_path, ppt_path)
        return jsonify({'message': 'Success', 'download_url': f'/download/{ppt_filename}'})
    except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/convert-to-word', methods=['POST'])
def convert_to_word():
    if 'file' not in request.files: return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    filename = secure_filename(file.filename)
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(pdf_path)
    word_filename = filename.rsplit('.', 1)[0] + '.docx'
    word_path = os.path.join(app.config['DOWNLOAD_FOLDER'], word_filename)
    try:
        cv = Converter(pdf_path)
        cv.convert(word_path, start=0, end=None)
        cv.close()
        return jsonify({'message': 'Success', 'download_url': f'/download/{word_filename}'})
    except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    return send_file(os.path.join(app.config['DOWNLOAD_FOLDER'], filename), as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)