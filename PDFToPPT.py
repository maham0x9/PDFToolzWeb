import os
import io  # New import for handling memory streams
import fitz  # PyMuPDF
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
from werkzeug.utils import secure_filename

app = Flask(__name__)
CORS(app)

# Configure folders
UPLOAD_FOLDER = 'uploads'
DOWNLOAD_FOLDER = 'downloads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DOWNLOAD_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['DOWNLOAD_FOLDER'] = DOWNLOAD_FOLDER

# --- HELPER: PDF TO PPTX FUNCTION (FIXED) ---
def convert_pdf_to_pptx_logic(pdf_path, pptx_path):
    # 1. Create a new Presentation
    prs = Presentation()
    
    # 2. Open PDF
    doc = fitz.open(pdf_path)
    
    # 3. Loop through pages
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # Render page to image (High Quality)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        
        # FIX: Save to Memory (RAM) instead of Disk to avoid Permission Errors
        img_stream = io.BytesIO(pix.tobytes("png"))
        
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add image from memory stream to slide
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            
    doc.close()
    prs.save(pptx_path)

# --- ROUTES ---

@app.route('/convert-to-word', methods=['POST'])
def convert_to_word():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        
        word_filename = filename.rsplit('.', 1)[0] + '.docx'
        word_path = os.path.join(app.config['DOWNLOAD_FOLDER'], word_filename)
        
        try:
            cv = Converter(pdf_path)
            cv.convert(word_path, start=0, end=None)
            cv.close()
            return jsonify({'message': 'Success', 'download_url': f'/download/{word_filename}'})
        except Exception as e:
            return jsonify({'error': str(e)}), 500

@app.route('/convert-to-ppt', methods=['POST'])
def convert_to_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)

        ppt_filename = filename.rsplit('.', 1)[0] + '.pptx'
        ppt_path = os.path.join(app.config['DOWNLOAD_FOLDER'], ppt_filename)

        try:
            convert_pdf_to_pptx_logic(pdf_path, ppt_path)

            return jsonify({
                'message': 'Conversion successful',
                'download_url': f'/download/{ppt_filename}'
            })
        except Exception as e:
            print(f"Error: {e}")
            return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    return send_file(os.path.join(app.config['DOWNLOAD_FOLDER'], filename), as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True, port=5000)