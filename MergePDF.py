import os
import io
import fitz  # PyMuPDF
import pdfplumber
import pandas as pd
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
from werkzeug.utils import secure_filename

app = Flask(__name__)
CORS(app)

# Configure folders
UPLOAD_FOLDER = 'uploads'
DOWNLOAD_FOLDER = 'downloads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DOWNLOAD_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['DOWNLOAD_FOLDER'] = DOWNLOAD_FOLDER

# --- HELPER FUNCTIONS ---

def convert_pdf_to_pptx_logic(pdf_path, pptx_path):
    prs = Presentation()
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_stream = io.BytesIO(pix.tobytes("png"))
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    doc.close()
    prs.save(pptx_path)

def convert_pdf_to_excel_logic(pdf_path, excel_path):
    with pdfplumber.open(pdf_path) as pdf:
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            tables_found = False
            for i, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                if tables:
                    tables_found = True
                    for j, table in enumerate(tables):
                        df = pd.DataFrame(table)
                        df = df.replace(r'\n', ' ', regex=True)
                        sheet_name = f"Page{i+1}_Table{j+1}"
                        df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
            if not tables_found:
                df = pd.DataFrame(["No detected tables in this PDF."])
                df.to_excel(writer, sheet_name="Info", index=False, header=False)

# --- ROUTES ---

@app.route('/convert-to-word', methods=['POST'])
def convert_to_word():
    # (Existing Word Logic - kept brief)
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        word_filename = filename.rsplit('.', 1)[0] + '.docx'
        word_path = os.path.join(app.config['DOWNLOAD_FOLDER'], word_filename)
        try:
            cv = Converter(pdf_path)
            cv.convert(word_path, start=0, end=None)
            cv.close()
            return jsonify({'message': 'Success', 'download_url': f'/download/{word_filename}'})
        except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/convert-to-ppt', methods=['POST'])
def convert_to_ppt():
    # (Existing PPT Logic - kept brief)
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        ppt_filename = filename.rsplit('.', 1)[0] + '.pptx'
        ppt_path = os.path.join(app.config['DOWNLOAD_FOLDER'], ppt_filename)
        try:
            convert_pdf_to_pptx_logic(pdf_path, ppt_path)
            return jsonify({'message': 'Success', 'download_url': f'/download/{ppt_filename}'})
        except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/convert-to-excel', methods=['POST'])
def convert_to_excel():
    # (Existing Excel Logic - kept brief)
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        excel_filename = filename.rsplit('.', 1)[0] + '.xlsx'
        excel_path = os.path.join(app.config['DOWNLOAD_FOLDER'], excel_filename)
        try:
            convert_pdf_to_excel_logic(pdf_path, excel_path)
            return jsonify({'message': 'Success', 'download_url': f'/download/{excel_filename}'})
        except Exception as e: return jsonify({'error': str(e)}), 500

# --- NEW: MERGE PDF ROUTE ---
@app.route('/merge-pdfs', methods=['POST'])
def merge_pdfs():
    # request.files.getlist pulls ALL files sent with the key 'files'
    uploaded_files = request.files.getlist('files')
    
    if not uploaded_files or uploaded_files[0].filename == '':
        return jsonify({'error': 'No files selected'}), 400

    try:
        # 1. Create a blank PDF to hold the result
        result_doc = fitz.open()

        # 2. Iterate through uploaded files
        for file in uploaded_files:
            # Read file stream directly into memory (no need to save temp files)
            file_stream = file.read()
            
            # Open the PDF from memory
            src_doc = fitz.open("pdf", file_stream)
            
            # Insert all pages from source to result
            result_doc.insert_pdf(src_doc)
            src_doc.close()

        # 3. Generate a meaningful filename (e.g., merged_timestamp.pdf or merged_firstfilename.pdf)
        first_name = secure_filename(uploaded_files[0].filename).rsplit('.', 1)[0]
        output_filename = f"Merged_{first_name}_and_others.pdf"
        output_path = os.path.join(app.config['DOWNLOAD_FOLDER'], output_filename)

        # 4. Save the result
        result_doc.save(output_path)
        result_doc.close()

        return jsonify({
            'message': 'Merge successful',
            'download_url': f'/download/{output_filename}'
        })

    except Exception as e:
        print(f"Merge Error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    return send_file(os.path.join(app.config['DOWNLOAD_FOLDER'], filename), as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True, port=5000)