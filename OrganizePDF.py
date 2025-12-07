import os
import io
import fitz  # PyMuPDF
import pdfplumber
import pandas as pd
import zipfile
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
from werkzeug.utils import secure_filename

app = Flask(__name__)
CORS(app)

# Configure folders
UPLOAD_FOLDER = 'uploads'
DOWNLOAD_FOLDER = 'downloads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DOWNLOAD_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['DOWNLOAD_FOLDER'] = DOWNLOAD_FOLDER

# --- HELPERS ---
def convert_pdf_to_pptx_logic(pdf_path, pptx_path):
    prs = Presentation()
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_stream = io.BytesIO(pix.tobytes("png"))
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    doc.close()
    prs.save(pptx_path)

def convert_pdf_to_excel_logic(pdf_path, excel_path):
    with pdfplumber.open(pdf_path) as pdf:
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            tables_found = False
            for i, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                if tables:
                    tables_found = True
                    for j, table in enumerate(tables):
                        df = pd.DataFrame(table)
                        df = df.replace(r'\n', ' ', regex=True)
                        sheet_name = f"Page{i+1}_Table{j+1}"
                        df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
            if not tables_found:
                df = pd.DataFrame(["No detected tables in this PDF."])
                df.to_excel(writer, sheet_name="Info", index=False, header=False)

def parse_page_string(order_str, total_pages):
    """Parses strings like '1, 3-5, 2' into a list of 0-based indices [0, 2, 3, 4, 1]"""
    selected_pages = []
    if not order_str:
        return list(range(total_pages)) # Default to original order
    
    parts = order_str.split(',')
    for part in parts:
        part = part.strip()
        if not part: continue
        try:
            if '-' in part:
                start, end = map(int, part.split('-'))
                # Handle "1-5" (forward) and "5-1" (reverse)
                step = 1 if start <= end else -1
                # Adjust to 0-based index
                r_start = max(1, start) - 1
                r_end = min(total_pages, end) - 1
                
                # Python range end is exclusive, so we add step
                # BUT be careful with reverse ranges in python
                if step == 1:
                    selected_pages.extend(range(r_start, r_end + 1))
                else:
                    selected_pages.extend(range(r_start, r_end - 1, -1))
            else:
                p = int(part) - 1
                if 0 <= p < total_pages:
                    selected_pages.append(p)
        except ValueError:
            pass # Ignore invalid chunks
    
    # If user entered garbage, fallback to all pages
    return selected_pages if selected_pages else list(range(total_pages))

# --- ROUTES ---

@app.route('/convert-to-word', methods=['POST'])
def convert_to_word():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        word_filename = filename.rsplit('.', 1)[0] + '.docx'
        word_path = os.path.join(app.config['DOWNLOAD_FOLDER'], word_filename)
        try:
            cv = Converter(pdf_path)
            cv.convert(word_path, start=0, end=None)
            cv.close()
            return jsonify({'message': 'Success', 'download_url': f'/download/{word_filename}'})
        except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/convert-to-ppt', methods=['POST'])
def convert_to_ppt():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        ppt_filename = filename.rsplit('.', 1)[0] + '.pptx'
        ppt_path = os.path.join(app.config['DOWNLOAD_FOLDER'], ppt_filename)
        try:
            convert_pdf_to_pptx_logic(pdf_path, ppt_path)
            return jsonify({'message': 'Success', 'download_url': f'/download/{ppt_filename}'})
        except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/convert-to-excel', methods=['POST'])
def convert_to_excel():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        excel_filename = filename.rsplit('.', 1)[0] + '.xlsx'
        excel_path = os.path.join(app.config['DOWNLOAD_FOLDER'], excel_filename)
        try:
            convert_pdf_to_excel_logic(pdf_path, excel_path)
            return jsonify({'message': 'Success', 'download_url': f'/download/{excel_filename}'})
        except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/merge-pdfs', methods=['POST'])
def merge_pdfs():
    uploaded_files = request.files.getlist('files')
    if not uploaded_files or uploaded_files[0].filename == '': return jsonify({'error': 'No files selected'}), 400
    try:
        result_doc = fitz.open()
        for file in uploaded_files:
            file_stream = file.read()
            src_doc = fitz.open("pdf", file_stream)
            result_doc.insert_pdf(src_doc)
            src_doc.close()
        first_name = secure_filename(uploaded_files[0].filename).rsplit('.', 1)[0]
        output_filename = f"Merged_{first_name}_and_others.pdf"
        output_path = os.path.join(app.config['DOWNLOAD_FOLDER'], output_filename)
        result_doc.save(output_path)
        result_doc.close()
        return jsonify({'message': 'Merge successful', 'download_url': f'/download/{output_filename}'})
    except Exception as e: return jsonify({'error': str(e)}), 500

@app.route('/split-pdf', methods=['POST'])
def split_pdf():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)
        start_page = request.form.get('start_page', type=int)
        end_page = request.form.get('end_page', type=int)
        base_name = filename.rsplit('.', 1)[0]
        zip_filename = f"{base_name}_split_files.zip"
        zip_path = os.path.join(app.config['DOWNLOAD_FOLDER'], zip_filename)
        try:
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            start_idx = (start_page - 1) if start_page else 0
            end_idx = end_page if end_page else total_pages
            if start_idx < 0: start_idx = 0
            if end_idx > total_pages: end_idx = total_pages
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for page_num in range(start_idx, end_idx):
                    new_doc = fitz.open()
                    new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
                    pdf_bytes = new_doc.tobytes()
                    new_doc.close()
                    zipf.writestr(f"{base_name}_page_{page_num + 1}.pdf", pdf_bytes)
            doc.close()
            return jsonify({'message': 'Split successful', 'download_url': f'/download/{zip_filename}'})
        except Exception as e: return jsonify({'error': str(e)}), 500

# --- NEW: ORGANIZE PDF ROUTE ---
@app.route('/organize-pdf', methods=['POST'])
def organize_pdf():
    if 'file' not in request.files: return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '': return jsonify({'error': 'No selected file'}), 400

    if file:
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pdf_path)

        # Get the order string from frontend (e.g. "1, 3, 2")
        page_order = request.form.get('page_order', '')

        base_name = filename.rsplit('.', 1)[0]
        output_filename = f"{base_name}_organized.pdf"
        output_path = os.path.join(app.config['DOWNLOAD_FOLDER'], output_filename)

        try:
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            
            # Parse the user string into a list of page indices
            selected_indices = parse_page_string(page_order, total_pages)
            
            # Create a new PDF based on that selection
            # We use select() to keep only specific pages in specific order
            doc.select(selected_indices)
            
            doc.save(output_path)
            doc.close()

            return jsonify({
                'message': 'Organize successful',
                'download_url': f'/download/{output_filename}'
            })

        except Exception as e:
            print(f"Organize Error: {e}")
            return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    return send_file(os.path.join(app.config['DOWNLOAD_FOLDER'], filename), as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True, port=5000)